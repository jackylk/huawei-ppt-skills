#!/usr/bin/env python3
"""Create a 16:9 PPTX from a folder of slide images.

Images are sorted lexicographically and placed full-bleed, one image per slide.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--images", required=True, help="Folder containing slide images.")
    parser.add_argument("--output", required=True, help="Output .pptx path.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    image_dir = Path(args.images).expanduser().resolve()
    output = Path(args.output).expanduser().resolve()
    images = sorted(
        [
            p
            for p in image_dir.iterdir()
            if p.suffix.lower() in {".png", ".jpg", ".jpeg"}
        ]
    )
    if not images:
        raise SystemExit(f"No slide images found in {image_dir}")

    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for image in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image), 0, 0, width=prs.slide_width, height=prs.slide_height
        )

    prs.save(output)
    print(output)
    print(f"slides={len(images)}")


if __name__ == "__main__":
    main()
